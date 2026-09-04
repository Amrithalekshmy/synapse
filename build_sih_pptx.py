"""
Create SYNAPSE_SIH2026.pptx — a brand-new file in the exact SIH2026 template format.
Starts from the official SIH2026-IDEA-Presentation-Format.pptx and fills every
slide with SYNAPSE content from the specifications folder.

Never modifies the existing SYNAPSE_final.pptx.
"""
import copy
import shutil
from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.enum.text import PP_ALIGN

# ── output file ──────────────────────────────────────────────────────────────
TEMPLATE = "specifications/SIH2026-IDEA-Presentation-Format.pptx"
OUTPUT   = "SYNAPSE_SIH2026.pptx"
shutil.copy(TEMPLATE, OUTPUT)
prs = Presentation(OUTPUT)

# ── colours (match the spec images) ─────────────────────────────────────────
BLUE    = RGBColor(0x00, 0x70, 0xC0)   # SIH template blue
DARK    = RGBColor(0x1F, 0x39, 0x64)   # dark navy (headings)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLACK   = RGBColor(0x00, 0x00, 0x00)
GREEN   = RGBColor(0x37, 0x86, 0x10)   # section header green (from slide2.jpeg)
AMBER   = RGBColor(0xFF, 0xC0, 0x00)   # bottom bar gold
LBLUE   = RGBColor(0xDD, 0xEB, 0xF7)   # light blue table header
DBLUE   = RGBColor(0x1F, 0x39, 0x64)   # dark blue table header text
LTEAL   = RGBColor(0xE2, 0xEF, 0xDA)   # light green background
LPURP   = RGBColor(0xE2, 0xD9, 0xF3)   # light purple
RED     = RGBColor(0xC0, 0x00, 0x00)

# ── slide dimensions ─────────────────────────────────────────────────────────
W  = prs.slide_width   # 12192000
H  = prs.slide_height  # 6858000

# content region (same left/width as title placeholder)
CL = 609600    # content left
CW = 10980400  # content width
CT = 1080000   # content top (just below team oval)
CB = 6300000   # content bottom (just above footer bar)
CH = CB - CT   # content height ≈ 5220000

# ── helpers ──────────────────────────────────────────────────────────────────

def tf_run(tf, text, size=Pt(11), bold=False, color=None, align=None, space_before=None, italic=False, para=None):
    """Add a paragraph+run to a text frame, return the run."""
    if para is None:
        para = tf.add_paragraph()
    if align:
        para.alignment = align
    if space_before:
        para.space_before = space_before
    run = para.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = "Arial"
    if color:
        run.font.color.rgb = color
    return run, para


def add_box(slide, left, top, width, height, text="", size=Pt(11), bold=False,
            color=BLACK, bg=None, align=PP_ALIGN.LEFT, wrap=True):
    """Add a text box and return its text frame."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    if text:
        run = p.add_run()
        run.text = text
        run.font.size = size
        run.font.bold = bold
        run.font.name = "Arial"
        run.font.color.rgb = color
    if bg:
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = bg
    return tf


def add_section_header(slide, top, number, title, height=360000, color=GREEN):
    """Numbered section header bar spanning full content width."""
    tf = add_box(slide, CL, top, CW, height,
                 text=f"{number}. {title}",
                 size=Pt(13), bold=True, color=WHITE, bg=color)
    return tf, top + height


def add_bullet_box(slide, left, top, width, height, items, size=Pt(10.5), color=BLACK, indent="• "):
    """Add a text box with bullet items."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = (indent + item) if indent else item
        run.font.size = size
        run.font.name = "Arial"
        run.font.color.rgb = color
    return tf


def tbl_cell(cell, text, bold=False, size=Pt(9.5), bg=None, color=BLACK, align=PP_ALIGN.LEFT):
    """Set a table cell text with formatting."""
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.name = "Arial"
    run.font.color.rgb = color
    if bg:
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        solidFill = copy.deepcopy(tc.get_or_add_tcPr())
        # Use direct XML for cell background
        from pptx.oxml import parse_xml
        from pptx.oxml.ns import nsmap
        hex_color = f"{bg[0]:02X}{bg[1]:02X}{bg[2]:02X}"
        solidFillXml = f'<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><a:srgbClr val="{hex_color}"/></a:solidFill>'
        fill_elem = parse_xml(solidFillXml)
        existing = tcPr.find(qn('a:solidFill'))
        if existing is not None:
            tcPr.remove(existing)
        noFill = tcPr.find(qn('a:noFill'))
        if noFill is not None:
            tcPr.remove(noFill)
        tcPr.insert(0, fill_elem)


def set_oval_text(slide, text):
    """Change the 'Your Team Name' oval text."""
    for shape in slide.shapes:
        if 'Oval' in shape.name and shape.has_text_frame:
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = text
            run.font.size = Pt(8)
            run.font.bold = True
            run.font.name = "Arial"
            break


def change_title(slide, new_title, size=Pt(32)):
    """Change the title placeholder text."""
    for shape in slide.shapes:
        if shape.name == 'Title 1' and shape.has_text_frame:
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = new_title
            run.font.size = size
            run.font.bold = True
            run.font.name = "Times New Roman"
            run.font.color.rgb = DARK
            break


def clear_content_textbox(slide):
    """Clear the big placeholder TextBox 8 (instructions text)."""
    for shape in slide.shapes:
        if shape.name == 'TextBox 8' and shape.has_text_frame:
            shape.text_frame.clear()
            break


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE PAGE
# ═════════════════════════════════════════════════════════════════════════════
s1 = prs.slides[0]
for shape in s1.shapes:
    if shape.name == 'TextBox 9' and shape.has_text_frame:
        tf = shape.text_frame
        tf.clear()
        fields = [
            ("Problem Statement ID – ", "SIH26122"),
            ("Problem Statement Title – ", "Context-Aware L5/L6 Schedule Reconciliation via Intelligent Field Progress Capture"),
            ("Theme – ", "Smart Automation"),
            ("PS Category – ", "Software"),
            ("Team ID – ", "[Team ID]"),
            ("Team Name – ", "Tachyons"),
        ]
        first = True
        for label, value in fields:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_before = Pt(4)
            r1 = p.add_run()
            r1.text = label
            r1.font.size = Pt(13)
            r1.font.bold = True
            r1.font.name = "Arial"
            r1.font.color.rgb = WHITE
            r2 = p.add_run()
            r2.text = value
            r2.font.size = Pt(13)
            r2.font.name = "Arial"
            r2.font.color.rgb = WHITE
        break

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — SYNAPSE (Idea Title)
# ═════════════════════════════════════════════════════════════════════════════
s2 = prs.slides[1]
clear_content_textbox(s2)
change_title(s2, "SYNAPSE — AI-Powered Field Progress to Schedule Reconciliation")
set_oval_text(s2, "Tachyons")

top = CT  # 1080000

# ── Section 1 ─────────────────────────────────────────────────────────────
_, top = add_section_header(s2, top, "1", "DETAILED EXPLANATION OF THE PROPOSED SOLUTION")
tb = s2.shapes.add_textbox(CL, top, CW, 820000)
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = ("SYNAPSE is an AI-powered system that converts unstructured field progress reports "
            "(DPRs, site diaries, discipline spreadsheets, supervisor voice/text inputs) into structured "
            "execution events and automatically links them to the correct L5/L6 schedule activities in "
            "Primavera P6. It extracts execution events, matches using a 7-layer contextual model, handles "
            "uncertainty through a human-in-the-loop clarification agent, and updates the schedule with a "
            "full audit trail and analytics.")
run.font.size = Pt(11)
run.font.name = "Arial"
top += 840000

# ── Section 2 ─────────────────────────────────────────────────────────────
_, top = add_section_header(s2, top, "2", "HOW IT ADDRESSES THE PROBLEM", color=BLUE)
bullets2 = [
    "Eliminates manual reconciliation of unstructured field reports with schedule activities.",
    "Reduces planner effort from 3–5 hours/day to less than 30 minutes/day.",
    "Minimises mis-linking and errors using multi-layer contextual matching and confidence thresholds.",
    "Provides near real-time progress visibility for monitoring, delay analysis and forecasting.",
    "Creates structured historical data for analytics, risk prediction and institutional memory.",
]
add_bullet_box(s2, CL, top, CW, 820000, bullets2, size=Pt(11))
top += 840000

# ── Section 3 ─────────────────────────────────────────────────────────────
_, top = add_section_header(s2, top, "3", "INNOVATION AND UNIQUENESS OF THE SOLUTION", color=DARK)
bullets3 = [
    "7-Layer Contextual Matching: Semantic similarity with identifier, discipline, location, WBS, temporal and dependency context for accurate activity mapping.",
    "Human-in-the-Loop Adaptive Agent: Decides to auto-link, ask clarification, escalate or flag new activity based on confidence; learns from planner feedback.",
    "End-to-End Integration: Data ingestion → schedule update → audit trail → analytics in a single workflow.",
    "Active Learning from Feedback: Planner decisions continuously improve matching accuracy and decision policy.",
    "Explainable & Safe: Confidence scores, reasons, audit trail and controlled updates ensure trust and reliability.",
]
add_bullet_box(s2, CL, top, CW, 900000, bullets3, size=Pt(11))
top += 930000

# ── Bottom summary bar ────────────────────────────────────────────────────
bar_h = CB - top
bar = s2.shapes.add_textbox(CL, top, CW, bar_h)
tf_bar = bar.text_frame
tf_bar.word_wrap = False
bar.fill.solid()
bar.fill.fore_color.rgb = AMBER
p = tf_bar.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r1 = p.add_run()
r1.text = "SYNAPSE — bridging the execution–schedule gap:   "
r1.font.size = Pt(11.5)
r1.font.bold = True
r1.font.name = "Arial"
r1.font.color.rgb = DARK
for label, val in [("3–5 hrs/day → <30 min", "  ▪  "), ("24–72 hrs lag → Near real-time", "  ▪  "), ("10–15% mis-link → <2%", "  ▪  "), ("Manual workflow → Automated + HITL", "")]:
    rx = p.add_run()
    rx.text = label
    rx.font.size = Pt(11)
    rx.font.bold = False
    rx.font.name = "Arial"
    rx.font.color.rgb = DARK
    if val:
        rv = p.add_run()
        rv.text = val
        rv.font.size = Pt(11)
        rv.font.name = "Arial"
        rv.font.color.rgb = DARK


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — TECHNICAL APPROACH
# ═════════════════════════════════════════════════════════════════════════════
s3 = prs.slides[2]
clear_content_textbox(s3)
change_title(s3, "TECHNICAL APPROACH")
set_oval_text(s3, "Tachyons")

top = CT

# ── Architecture header ───────────────────────────────────────────────────
_, top = add_section_header(s3, top, "1", "SYSTEM ARCHITECTURE — 8-Stage Pipeline", color=DARK)

# 4 architecture boxes on the left column (stages 1-4)
arch_col_w = CW // 2 - 30000
arch_stages_left = [
    ("1. INPUT SOURCES",
     "DPRs (PDF/Excel/Text) · Site Diaries / Supervisor Notes · Discipline Spreadsheets · Messages/Emails · Other Project Documents"),
    ("2. INGESTION & PREPROCESSING",
     "Document Ingestion → Data Cleaning & Normalisation → OCR (if needed) → Text Segmentation & Chunking → Metadata Extraction (Date, Discipline, Location)"),
    ("3. EVENT EXTRACTION (What happened?)",
     "LLM/NLP Model + Rule-based fallback → Extracted Events: Activities, Quantities, Location, Resources, Dates, Status\nExample: 'Excavation completed · Rebar fixing done · Concrete poured'"),
    ("4. CONTEXTUAL ACTIVITY MATCHING (Which L5/L6?)",
     "7-Layer Matching: ① Semantic Similarity ② Activity Identifier ③ Discipline ④ Location/Zone ⑤ WBS/Hierarchy ⑥ Temporal Context ⑦ Dependency → Best Match + Confidence Score"),
]
arch_stages_right = [
    ("5. DECISION & HUMAN-IN-THE-LOOP",
     "HIGH confidence (≥ 0.85) → Auto-link\nMEDIUM confidence → Ask Clarification / Human Review\nLOW confidence (< 0.65) → Escalate to Planner / Create New Activity"),
    ("6. SCHEDULE UPDATE & AUDIT TRAIL",
     "Update Primavera P6 Actuals (% Complete, Status, Dates) · Audit Trail: Who, What, When, Confidence, Reason · Versioning & Change History"),
    ("7. ANALYTICS & REPORTING",
     "Progress Dashboards (Real-time Visibility) · Delay / Variance Analysis · Risk Indicators · Institutional Memory (Knowledge Base)"),
    ("8. FEEDBACK & LEARNING LOOP",
     "Human Feedback → Store Labelled Pairs → Active Learning (Select Uncertain Cases) → Model/Policy Update → Improved Matching & Decision Policy"),
]

box_h = 530000
gap = 20000
y_left = top
y_right = top

for title, desc in arch_stages_left:
    # title bar
    hdr = s3.shapes.add_textbox(CL, y_left, arch_col_w, 240000)
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = BLUE
    hf = hdr.text_frame
    p = hf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(9.5)
    r.font.bold = True
    r.font.name = "Arial"
    r.font.color.rgb = WHITE
    hf.word_wrap = True
    # content
    cnt = s3.shapes.add_textbox(CL, y_left + 240000, arch_col_w, box_h - 240000)
    cnt.fill.solid()
    cnt.fill.fore_color.rgb = LBLUE
    cf = cnt.text_frame
    cf.word_wrap = True
    p = cf.paragraphs[0]
    r = p.add_run()
    r.text = desc
    r.font.size = Pt(9)
    r.font.name = "Arial"
    r.font.color.rgb = BLACK
    y_left += box_h + gap

right_start = CL + arch_col_w + 60000
for title, desc in arch_stages_right:
    hdr = s3.shapes.add_textbox(right_start, y_right, arch_col_w, 240000)
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = GREEN
    hf = hdr.text_frame
    p = hf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(9.5)
    r.font.bold = True
    r.font.name = "Arial"
    r.font.color.rgb = WHITE
    hf.word_wrap = True
    cnt = s3.shapes.add_textbox(right_start, y_right + 240000, arch_col_w, box_h - 240000)
    cnt.fill.solid()
    cnt.fill.fore_color.rgb = LTEAL
    cf = cnt.text_frame
    cf.word_wrap = True
    p = cf.paragraphs[0]
    r = p.add_run()
    r.text = desc
    r.font.size = Pt(9)
    r.font.name = "Arial"
    r.font.color.rgb = BLACK
    y_right += box_h + gap

top = max(y_left, y_right) + 40000

# ── Tech stack table ──────────────────────────────────────────────────────
_, top = add_section_header(s3, top, "2", "TECHNOLOGY STACK", color=DARK, height=320000)
tbl_top = top
tbl_rows_data = [
    ("Layer", "Technology", "Purpose"),
    ("User Interface", "HTML + CSS + Vanilla JavaScript", "Planner dashboard, review screen, clarification/approval"),
    ("Backend / API", "Python + FastAPI", "REST APIs, workflow orchestration, communication between modules"),
    ("Data Validation", "Pydantic", "Structured validation of extracted ExecutionEvent data"),
    ("Document Processing", "pandas + pdfplumber", "Processing DPRs, discipline spreadsheets and PDF reports"),
    ("Event Extraction", "LLM/NLP + Rule-based fallback", "Extract activity, location, discipline, dates, status"),
    ("Semantic Matching", "Sentence-Transformers / MiniLM", "Semantic similarity between field events and L5/L6 activities"),
    ("Contextual Matching", "Python scoring engine (7-layer)", "Combine semantic, identifier, discipline, WBS, temporal, dependency"),
    ("Decision Engine", "Confidence-based Router", "Auto-link, ask clarification, escalate or flag new activity"),
]
tbl_h = CB - tbl_top
tbl = s3.shapes.add_table(len(tbl_rows_data), 3, CL, tbl_top, CW, tbl_h).table
col_widths = [int(CW * 0.20), int(CW * 0.33), int(CW * 0.47)]
for i, w in enumerate(col_widths):
    tbl.columns[i].width = w

for ri, row_data in enumerate(tbl_rows_data):
    is_hdr = ri == 0
    bg = DBLUE if is_hdr else (LBLUE if ri % 2 == 0 else WHITE)
    txt_color = WHITE if is_hdr else BLACK
    for ci, cell_text in enumerate(row_data):
        tbl_cell(tbl.rows[ri].cells[ci], cell_text,
                 bold=is_hdr or ci == 0,
                 size=Pt(9) if not is_hdr else Pt(9.5),
                 bg=bg, color=txt_color)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — FEASIBILITY AND VIABILITY
# ═════════════════════════════════════════════════════════════════════════════
s4 = prs.slides[3]
clear_content_textbox(s4)
change_title(s4, "FEASIBILITY AND VIABILITY")
set_oval_text(s4, "Tachyons")

top = CT
col1_w = int(CW * 0.295)
col2_w = int(CW * 0.295)
col3_w = CW - col1_w - col2_w - 60000
col2_l = CL + col1_w + 30000
col3_l = col2_l + col2_w + 30000

# column headers
for col_l, col_w, label, col_color in [
    (CL,    col1_w, "1. TECHNICAL FEASIBILITY",  GREEN),
    (col2_l, col2_w, "2. OPERATIONAL VIABILITY",  BLUE),
    (col3_l, col3_w, "3. KEY RISKS & MITIGATION", DARK),
]:
    hdr = s4.shapes.add_textbox(col_l, top, col_w, 380000)
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = col_color
    hf = hdr.text_frame
    hf.word_wrap = True
    p = hf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.name = "Arial"
    r.font.color.rgb = WHITE

col_content_top = top + 400000
col_h = CB - col_content_top - 500000  # leave room for footer banner

# Col 1: Technical Feasibility bullets
feasibility = [
    "CPU-based MiniLM semantic matching; no GPU required for the MVP.",
    "Python + FastAPI + Pydantic backend; pandas / pdfplumber for DPR and spreadsheet processing.",
    "P6 schedule data consumed through controlled export/import; production can use P6 API where available.",
    "Prototype demonstrated on a 42-activity schedule with 3 days of daily progress reports and 49 historical records.",
    "100% open-source stack — zero proprietary scheduling-tool licences.",
    "Works fully offline: rule-based extraction + TF-IDF matching fallback active at all times.",
]
bx = s4.shapes.add_textbox(CL, col_content_top, col1_w, col_h)
bx.fill.solid()
bx.fill.fore_color.rgb = LTEAL
btf = bx.text_frame
btf.word_wrap = True
first = True
for item in feasibility:
    p = btf.paragraphs[0] if first else btf.add_paragraph()
    first = False
    p.space_before = Pt(4)
    r = p.add_run()
    r.text = "• " + item
    r.font.size = Pt(10)
    r.font.name = "Arial"
    r.font.color.rgb = BLACK

# Col 2: Operational Viability bullets
viability = [
    "Fits the existing planner workflow: AI assists; planner retains control over uncertain cases.",
    "Confidence routing: high → auto-link; medium → clarification/HITL; low → planner review / new activity.",
    "Every approved update records confidence, reason, user and timestamp for full auditability.",
    "Scales from the prototype to large L5/L6 schedules with candidate retrieval and vector search.",
    "Zero workflow change — supervisors submit the same DPRs they already write.",
    "Voice input via browser — supervisor speaks, SYNAPSE understands (Web Speech API).",
]
bx2 = s4.shapes.add_textbox(col2_l, col_content_top, col2_w, col_h)
bx2.fill.solid()
bx2.fill.fore_color.rgb = LBLUE
btf2 = bx2.text_frame
btf2.word_wrap = True
first = True
for item in viability:
    p = btf2.paragraphs[0] if first else btf2.add_paragraph()
    first = False
    p.space_before = Pt(4)
    r = p.add_run()
    r.text = "• " + item
    r.font.size = Pt(10)
    r.font.name = "Arial"
    r.font.color.rgb = BLACK

# Col 3: Risks table
risk_rows = [
    ("Risk", "Impact", "Mitigation"),
    ("Ambiguous field language", "High", "7-layer contextual matching + targeted clarification"),
    ("Wrong auto-link", "High", "Conservative threshold + HITL + audit trail + rollback"),
    ("Domain shift / limited labelled data", "Medium", "Active learning from planner feedback"),
    ("No live OIL data", "High", "Synthetic dataset (42 activities, 3 DPRs, 25 labelled pairs)"),
    ("Planner adoption", "Medium", "Augments, not replaces, planner decisions"),
]
risk_tbl_h = col_h
risk_tbl = s4.shapes.add_table(len(risk_rows), 3, col3_l, col_content_top, col3_w, risk_tbl_h).table
risk_col_w = [int(col3_w * 0.38), int(col3_w * 0.16), int(col3_w * 0.46)]
for i, w in enumerate(risk_col_w):
    risk_tbl.columns[i].width = w
for ri, row_data in enumerate(risk_rows):
    is_hdr = ri == 0
    bg = DBLUE if is_hdr else (RGBColor(0xFF, 0xEB, 0xEB) if ri % 2 == 1 else WHITE)
    txt_color = WHITE if is_hdr else BLACK
    for ci, text in enumerate(row_data):
        tbl_cell(risk_tbl.rows[ri].cells[ci], text,
                 bold=is_hdr or ci == 0,
                 size=Pt(8.5) if not is_hdr else Pt(9),
                 bg=bg, color=txt_color)

# Footer banner
banner_top = col_content_top + col_h + 60000
banner = s4.shapes.add_textbox(CL, banner_top, CW, CB - banner_top)
banner.fill.solid()
banner.fill.fore_color.rgb = AMBER
btf = banner.text_frame
p = btf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "FEASIBLE MVP TODAY    •    SCALABLE PRODUCTION PATH    •    HUMAN-CONTROLLED AI"
r.font.size = Pt(13)
r.font.bold = True
r.font.name = "Arial"
r.font.color.rgb = DARK


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — IMPACT AND BENEFITS
# ═════════════════════════════════════════════════════════════════════════════
s5 = prs.slides[4]
clear_content_textbox(s5)
change_title(s5, "IMPACT AND BENEFITS")
set_oval_text(s5, "Tachyons")

top = CT
# 3 columns: target audience, benefits, SDG
col1_w = int(CW * 0.27)
col2_w = int(CW * 0.40)
col3_w = CW - col1_w - col2_w - 60000
col2_l = CL + col1_w + 30000
col3_l = col2_l + col2_w + 30000

for col_l, col_w, label, col_color in [
    (CL,    col1_w, "1. POTENTIAL IMPACT ON TARGET AUDIENCE", GREEN),
    (col2_l, col2_w, "2. BENEFITS OF THE SOLUTION",            BLUE),
    (col3_l, col3_w, "3. SDG ALIGNMENT",                       DARK),
]:
    hdr = s5.shapes.add_textbox(col_l, top, col_w, 380000)
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = col_color
    hf = hdr.text_frame
    hf.word_wrap = True
    p = hf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.name = "Arial"
    r.font.color.rgb = WHITE

col_content_top = top + 400000
col_h = 4150000

# Col 1: Target audience
audience = [
    ("Project Planners / Engineers",
     "3–5× faster reconciliation; more time for planning and decision-making; reduced manual effort."),
    ("Management",
     "Real-time visibility into project execution; earlier delay detection and better control."),
    ("Field Teams",
     "Less back-and-forth for clarifications; clearer communication and faster approvals."),
    ("Organisation",
     "Improved data quality, auditability and institutional memory across projects."),
    ("End Clients / Stakeholders",
     "Projects delivered on schedule with better transparency and reliability."),
]
bx = s5.shapes.add_textbox(CL, col_content_top, col1_w, col_h)
bx.fill.solid()
bx.fill.fore_color.rgb = LTEAL
btf = bx.text_frame
btf.word_wrap = True
first = True
for group, desc in audience:
    p = btf.paragraphs[0] if first else btf.add_paragraph()
    first = False
    p.space_before = Pt(5)
    r1 = p.add_run()
    r1.text = group + "\n"
    r1.font.size = Pt(10)
    r1.font.bold = True
    r1.font.name = "Arial"
    r1.font.color.rgb = DARK
    r2 = p.add_run()
    r2.text = desc
    r2.font.size = Pt(9.5)
    r2.font.name = "Arial"
    r2.font.color.rgb = BLACK

# Col 2: Benefits table
benefit_rows = [
    ("Benefit", "Impact"),
    ("Time Savings", "Reduces planner effort from 3–5 hrs/day to <30 minutes/day"),
    ("Higher Accuracy", "Multi-layer matching + confidence routing minimises mis-linking"),
    ("Real-time Visibility", "Near real-time updates enable proactive delay & risk identification"),
    ("Better Decisions", "Data-driven insights for planning, resource allocation & forecasting"),
    ("Data Quality & Auditability", "Every update captured with confidence, reason & user"),
    ("Scalability", "Works for small projects; scales to large L5/L6 schedules"),
    ("Cost Efficiency", "Reduces rework, delays and escalations → lower project costs"),
    ("Knowledge Retention", "Creates structured historical data; learns from planner feedback"),
]
benefit_tbl = s5.shapes.add_table(len(benefit_rows), 2, col2_l, col_content_top, col2_w, col_h).table
bc_w = [int(col2_w * 0.35), int(col2_w * 0.65)]
for i, w in enumerate(bc_w):
    benefit_tbl.columns[i].width = w
for ri, row_data in enumerate(benefit_rows):
    is_hdr = ri == 0
    bg = DBLUE if is_hdr else (LBLUE if ri % 2 == 0 else WHITE)
    txt_color = WHITE if is_hdr else BLACK
    for ci, text in enumerate(row_data):
        tbl_cell(benefit_tbl.rows[ri].cells[ci], text,
                 bold=is_hdr or ci == 0,
                 size=Pt(9) if not is_hdr else Pt(9.5),
                 bg=bg, color=txt_color)

# Col 3: SDG alignment
sdgs = [
    ("SDG 9", "Industry, Innovation & Infrastructure",
     "Resilient infrastructure through AI-powered project management and digital transformation."),
    ("SDG 8", "Decent Work & Economic Growth",
     "Improves productivity, reduces manual effort and delays, leading to efficient use of resources."),
    ("SDG 12", "Responsible Consumption & Production",
     "Optimises resource utilisation, reduces waste and promotes efficient project execution."),
    ("SDG 13", "Climate Action",
     "Better planning & monitoring, helping reduce resource and carbon footprint."),
    ("SDG 17", "Partnerships for the Goals",
     "Strengthens collaboration between field teams, planners and stakeholders through data-driven workflows."),
]
bx3 = s5.shapes.add_textbox(col3_l, col_content_top, col3_w, col_h)
bx3.fill.solid()
bx3.fill.fore_color.rgb = LPURP
btf3 = bx3.text_frame
btf3.word_wrap = True
first = True
for sdg_num, sdg_title, sdg_desc in sdgs:
    p = btf3.paragraphs[0] if first else btf3.add_paragraph()
    first = False
    p.space_before = Pt(5)
    r1 = p.add_run()
    r1.text = f"{sdg_num}: {sdg_title}\n"
    r1.font.size = Pt(10)
    r1.font.bold = True
    r1.font.name = "Arial"
    r1.font.color.rgb = DARK
    r2 = p.add_run()
    r2.text = sdg_desc
    r2.font.size = Pt(9)
    r2.font.name = "Arial"
    r2.font.color.rgb = BLACK

# Overall impact banner
banner_top = col_content_top + col_h + 50000
banner = s5.shapes.add_textbox(CL, banner_top, CW, CB - banner_top)
banner.fill.solid()
banner.fill.fore_color.rgb = AMBER
btf = banner.text_frame
btf.word_wrap = True
p = btf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r1 = p.add_run()
r1.text = "OVERALL IMPACT: "
r1.font.size = Pt(11)
r1.font.bold = True
r1.font.name = "Arial"
r1.font.color.rgb = DARK
r2 = p.add_run()
r2.text = ("SYNAPSE transforms unstructured field reports into trusted, real-time execution intelligence — "
           "enabling smarter decisions, on-time delivery and sustainable infrastructure development.")
r2.font.size = Pt(11)
r2.font.name = "Arial"
r2.font.color.rgb = DARK


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — RESEARCH AND REFERENCES
# ═════════════════════════════════════════════════════════════════════════════
s6 = prs.slides[5]
clear_content_textbox(s6)
change_title(s6, "RESEARCH AND REFERENCES")
set_oval_text(s6, "Tachyons")

top = CT

# ── What our research establishes ────────────────────────────────────────
_, top = add_section_header(s6, top, "1", "WHAT OUR RESEARCH ESTABLISHES", color=DARK)

research_items = [
    ("Problem is real — ",
     "EPC planners spend 3–5 hrs/day manually reconciling site reports against schedules with 500–5,000 "
     "activities; updates reach the schedule 24–72 hrs late."),
    ("Prior tools fail — ",
     "Primavera P6 and MS Project capture site progress only through manual entry. "
     "No NLP, no auto-linking, no learning."),
    ("The gap — ",
     "No existing EPC tool reads free-text reports, auto-links them to L5/L6 activities, clarifies ambiguity, "
     "detects cross-source conflicts, or learns from corrections."),
    ("Why our method works — ",
     "Sentence-BERT semantic matching (Reimers & Gurevych, 2019), contextual bandits (Li et al., 2010), "
     "and BFS graph analysis are individually established — SYNAPSE combines them for EPC reconciliation."),
    ("Validated — ",
     "Synthetic Oil India dataset: 42 L5/L6 activities, 5 disciplines, 3 DPRs, "
     "49 historical records, 25 labelled match pairs. 220 automated tests passing."),
]

res_box = s6.shapes.add_textbox(CL, top, CW, 2200000)
res_tf = res_box.text_frame
res_tf.word_wrap = True
first = True
for bold_text, body_text in research_items:
    p = res_tf.paragraphs[0] if first else res_tf.add_paragraph()
    first = False
    p.space_before = Pt(5)
    r1 = p.add_run()
    r1.text = "• " + bold_text
    r1.font.size = Pt(11)
    r1.font.bold = True
    r1.font.name = "Arial"
    r1.font.color.rgb = DARK
    r2 = p.add_run()
    r2.text = body_text
    r2.font.size = Pt(11)
    r2.font.name = "Arial"
    r2.font.color.rgb = BLACK
top += 2250000

# ── References ────────────────────────────────────────────────────────────
_, top = add_section_header(s6, top, "2", "REFERENCES — AUTHORITATIVE SOURCES ONLY", color=BLUE)
ref_box = s6.shapes.add_textbox(CL, top, CW, CB - top)
ref_tf = ref_box.text_frame
ref_tf.word_wrap = True

ref_sections = [
    ("PROBLEM EVIDENCE", True, None),
    ("Oil India Limited — Annual Reports & EPC Project Documentation", False, None),
    ("EPC Schedule Management Industry Reports (manual reconciliation workflow)", False, None),
    ("", False, None),
    ("STATE OF THE ART", True, None),
    ("Reimers & Gurevych (2019), Sentence-BERT: Sentence Embeddings using Siamese Networks — EMNLP 2019", False, None),
    ("Settles (2009), Active Learning Literature Survey — University of Wisconsin–Madison, TR-1648", False, None),
    ("Li, Chu, Langford & Schapire (2010), A Contextual-Bandit Approach to Personalized Recommendation — WWW", False, None),
    ("Oracle Corporation — Primavera P6 EPPM Documentation (manual data-entry workflow)", False, None),
    ("", False, None),
    ("IMPLEMENTATION", True, None),
    ("Hugging Face — sentence-transformers/all-MiniLM-L6-v2 (384-dim, Apache 2.0)", False, None),
    ("Oracle — Primavera P6 XER Export Format Specification", False, None),
    ("github.com/Amrithalekshmy/synapse — Full source code + synthetic dataset (open-source)", False, None),
]
first = True
for text, is_hdr, _ in ref_sections:
    p = ref_tf.paragraphs[0] if first else ref_tf.add_paragraph()
    first = False
    if not text:
        p.add_run().text = ""
        continue
    p.space_before = Pt(2 if not is_hdr else 6)
    r = p.add_run()
    r.text = ("▶ " if is_hdr else "  • ") + text
    r.font.size = Pt(10.5) if is_hdr else Pt(10)
    r.font.bold = is_hdr
    r.font.name = "Arial"
    r.font.color.rgb = BLUE if is_hdr else BLACK


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Remove instructions slide (delete from presentation)
# We keep it blank to avoid changing slide count issues; just clear it.
# ═════════════════════════════════════════════════════════════════════════════
# (Leave as-is — judges won't be confused; it's the official SIH instructions page)

# ── save ─────────────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f"✓ {OUTPUT} created successfully.")
print("  Slide 1: Title page — SYNAPSE / Tachyons / SIH26122")
print("  Slide 2: Idea Title — 3 sections (solution, problem, innovation) + metrics bar")
print("  Slide 3: Technical Approach — 8-box pipeline + 9-row tech stack table")
print("  Slide 4: Feasibility & Viability — 2 bullet columns + risks table + footer")
print("  Slide 5: Impact & Benefits — audience, benefits table, SDG + impact banner")
print("  Slide 6: Research & References — 5 findings + 3-section bibliography")
